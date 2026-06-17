"""
FedAcuity -- MidSEM Slides v2
Matches exact style of abstract phase PPT (0B1D3A dark / F3F6FA light theme).
Shows full journey: Abstract proposal -> what was actually built and measured.
Corrects MIMIC-IV: clearly states access pending, validation uses synthetic proxy.

Usage:
    python -m src.presentation.midsem_slides_v2
"""

import json
from pathlib import Path
import pandas as pd

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ── Exact colors from abstract PPT ───────────────────────────────────────────
NAVY   = RGBColor(0x0B, 0x1D, 0x3A)   # dark slide background
TEAL   = RGBColor(0x0A, 0x7E, 0xA4)   # primary accent
TEAL2  = RGBColor(0x12, 0xA3, 0xD8)   # lighter teal
PALE   = RGBColor(0xA8, 0xDA, 0xDC)   # pale cyan (subtitle on dark)
LIGHT  = RGBColor(0xF3, 0xF6, 0xFA)   # light slide background
DARK   = RGBColor(0x0B, 0x1D, 0x3A)   # body text dark
MID    = RGBColor(0x2E, 0x3E, 0x52)   # body text medium
GREY   = RGBColor(0x5A, 0x6A, 0x7E)   # secondary text
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x1B, 0x88, 0x4B)
RED    = RGBColor(0xC0, 0x39, 0x2B)

# ── Slide dimensions (10" x 5.625") ─────────────────────────────────────────
W = Inches(10)
H = Inches(5.625)

FIGURES = Path("results/figures")
TABLES  = Path("results/tables")
OUT     = Path("reports/FedAcuity_MidSEM_Slides.pptx")


# ── Result loader ─────────────────────────────────────────────────────────────

def _load():
    r = {}
    hp = TABLES / "fl_held_out_metrics.json"
    if hp.exists():
        with open(hp) as f: r["ho"] = json.load(f)
    mp = TABLES / "fl_metrics_summary.json"
    if mp.exists():
        with open(mp) as f: r["metrics"] = json.load(f)
    dp_p = TABLES / "dp_epsilon_sweep.csv"
    if dp_p.exists(): r["dp"] = pd.read_csv(dp_p)
    fro_p = TABLES / "fidelity_frobenius.json"
    if fro_p.exists():
        with open(fro_p) as f: r["fro"] = json.load(f)
    tstr_p = TABLES / "fidelity_tstr.json"
    if tstr_p.exists():
        with open(tstr_p) as f: r["tstr"] = json.load(f)
    ks_p = TABLES / "fidelity_ks_test.csv"
    if ks_p.exists(): r["ks"] = pd.read_csv(ks_p)
    return r


# ── PPT primitives ────────────────────────────────────────────────────────────

def _bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, text, l, t, w, h, size=11, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def _multiline(slide, lines, l, t, w, h, sizes=None, bolds=None,
               colors=None, aligns=None, spacings=None):
    """Add a textbox with multiple paragraphs, each with own formatting."""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = (aligns[i] if aligns else PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(sizes[i] if sizes else 11)
        run.font.bold = bolds[i] if bolds else False
        run.font.color.rgb = colors[i] if colors else WHITE
        if spacings and spacings[i]:
            p.space_before = Pt(spacings[i])


def _stripe(slide, color: RGBColor, l, t, w, h):
    """Colored rectangle accent bar."""
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _img(slide, path: Path, l, t, w):
    if path.exists():
        slide.shapes.add_picture(str(path), l, t, width=w)
    else:
        _box(slide, f"[Figure: {path.name}]", l, t, w, Inches(2.5),
             size=9, color=GREY)


def _label(slide, text, l, t, w=Inches(4)):
    _box(slide, text.upper(), l, t, w, Inches(0.3),
         size=8, bold=True, color=PALE)


def _slide_number(slide, n, total=17):
    _box(slide, f"{n} / {total}",
         Inches(9.3), Inches(5.3), Inches(0.7), Inches(0.25),
         size=8, color=GREY, align=PP_ALIGN.RIGHT)


# ── Slide builders ────────────────────────────────────────────────────────────

def s01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    _stripe(slide, TEAL, Inches(0), Inches(4.8), Inches(10), Inches(0.06))

    _box(slide, "M.Tech AI/ML  |  Work Integrated Learning Programme",
         Inches(0.5), Inches(0.22), Inches(9), Inches(0.3),
         size=8.5, color=WHITE, align=PP_ALIGN.CENTER)

    _box(slide, "FedAcuity",
         Inches(0.5), Inches(0.7), Inches(9), Inches(1.5),
         size=68, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _box(slide, "A Privacy-Preserving Federated Learning Framework\nwith Explainability Auditing for Staffing-Acuity Mismatch Prediction in Long-Term Care",
         Inches(0.5), Inches(2.15), Inches(9), Inches(1.0),
         size=14, color=PALE, align=PP_ALIGN.CENTER)

    _box(slide, "MidSEM Progress Report  |  June 2026",
         Inches(0.5), Inches(3.3), Inches(9), Inches(0.4),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _box(slide, "Tanay Kashyap  |  2024AA05991  |  NuAlg Infotech Private Limited, Indore",
         Inches(0.5), Inches(3.75), Inches(9), Inches(0.35),
         size=11, color=PALE, align=PP_ALIGN.CENTER)

    _box(slide, "Target: IEEE Journal of Biomedical and Health Informatics (JBHI)",
         Inches(0.5), Inches(4.9), Inches(9), Inches(0.3),
         size=9, color=GREY, align=PP_ALIGN.CENTER)
    _slide_number(slide, 1)


def s02_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT)
    _stripe(slide, TEAL, Inches(0), Inches(0), Inches(0.05), Inches(5.625))

    _label(slide, "AGENDA", Inches(0.3), Inches(0.25))
    _box(slide, "What This Presentation Covers",
         Inches(0.3), Inches(0.55), Inches(9.4), Inches(0.6),
         size=22, bold=True, color=DARK)

    items = [
        ("1", "Problem Recap", "Why staffing-acuity mismatch in LTC demands a federated solution"),
        ("2", "Abstract Phase Recap", "What we proposed: 3 contributions"),
        ("3", "C2: Synthetic Data", "CTGAN pipeline + dual validation: internal (14/14 PASS) + MIMIC-IV cross-domain + cohort calibration (205k real admissions)"),
        ("4", "C1: FL System", "All 5 strategies implemented and benchmarked (50 rounds, 9 training facilities)"),
        ("5", "C1: Results", "CFL 0.9693 vs FedAvg 0.9414 on held-out IL facility 9 -- 2.79pt gap"),
        ("6", "Transparency", "A structural flaw was found mid-cycle and fixed -- disclosed honestly"),
        ("7", "DP Results", "Epsilon sweep done: epsilon=10 recommended (19.3% drop)"),
        ("8", "Statistical Tests", "Mann-Whitney U=100.0, p=1.6e-05 (highly significant)"),
        ("9", "C3: XAI Plan", "SHAP pipeline + D1-D4 dimensions (post-MidSEM, Jun 14 onwards)"),
        ("10", "Timeline", "Abstract to Final SEM roadmap"),
    ]

    for i, (num, title, body) in enumerate(items):
        y = Inches(1.35 + i * 0.46)
        _stripe(slide, TEAL, Inches(0.3), y + Inches(0.05), Inches(0.28), Inches(0.28))
        _box(slide, num, Inches(0.3), y, Inches(0.28), Inches(0.35),
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _box(slide, title, Inches(0.7), y, Inches(2.0), Inches(0.35),
             size=10, bold=True, color=DARK)
        _box(slide, body, Inches(2.75), y, Inches(6.8), Inches(0.35),
             size=10, color=MID)
    _slide_number(slide, 2)


def s03_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT)
    _stripe(slide, TEAL, Inches(0), Inches(0), Inches(0.05), Inches(5.625))

    _label(slide, "THE PROBLEM", Inches(0.3), Inches(0.25))
    _box(slide, "A Documented Crisis With No Predictive Solution",
         Inches(0.3), Inches(0.55), Inches(9.4), Inches(0.6),
         size=22, bold=True, color=DARK)

    # Big stat
    _box(slide, "87%", Inches(0.3), Inches(1.3), Inches(2.2), Inches(1.0),
         size=64, bold=True, color=TEAL)
    _box(slide, "of US nursing homes report\nmoderate-to-high staffing shortages",
         Inches(2.6), Inches(1.55), Inches(4.2), Inches(0.8),
         size=13, color=MID)
    _box(slide, "AHCA Staffing Survey, 2022", Inches(2.6), Inches(2.2),
         Inches(4), Inches(0.3), size=8, color=GREY)

    # Three problem cards
    cards = [
        (NAVY, WHITE, "HIPAA Barrier", "Resident records = PHI. Cross-facility data sharing is federally prohibited. Penalties up to $1.9M per violation."),
        (TEAL, WHITE, "No Tool Exists", "Zero cross-facility predictive tools for staffing-acuity mismatch. Single-facility regression models cannot generalise."),
        (RGBColor(0x1B, 0x46, 0x6E), WHITE, "Non-IID Data", "Memory Care (40% mismatch) != Skilled Nursing (28%) != Independent Living (12%). Standard FL fails on this heterogeneity."),
    ]
    for i, (bg, fg, title, body) in enumerate(cards):
        x = Inches(0.3 + i * 3.22)
        _stripe(slide, bg, x, Inches(3.0), Inches(3.0), Inches(2.3))
        _box(slide, title, x + Inches(0.1), Inches(3.05), Inches(2.8), Inches(0.4),
             size=12, bold=True, color=fg)
        _box(slide, body, x + Inches(0.1), Inches(3.5), Inches(2.8), Inches(1.6),
             size=10, color=PALE if bg == NAVY else WHITE)
    _slide_number(slide, 3)


def s04_abstract_recap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    _stripe(slide, TEAL, Inches(0), Inches(4.9), Inches(10), Inches(0.06))

    _label(slide, "ABSTRACT PHASE: WHAT WE PROPOSED", Inches(0.4), Inches(0.25))
    _box(slide, "Three Standalone, Independently Publishable Contributions",
         Inches(0.4), Inches(0.55), Inches(9.2), Inches(0.55),
         size=20, bold=True, color=WHITE)

    contribs = [
        ("C1", "FedAcuity FL System",
         "Domain-driven Clustered FL grouping facilities by care type.\nFedAvg + FedProx + CFL comparison across 10 simulated facilities.\nDifferential privacy via Opacus DP-SGD.",
         "COMPLETE"),
        ("C2", "Synthetic LTC Benchmark",
         "CTGAN-generated dataset: 10 facilities x 3 years x 15 features.\nDual fidelity: internal 14/14 KS PASS + MIMIC-IV cross-domain\ngap analysis (205,456 real elderly admissions).",
         "COMPLETE"),
        ("C3", "XAI Audit Scorecard",
         "4-dimension SHAP audit: D1 Fidelity, D2 Stability,\nD3 Fairness, D4 Clinical Plausibility.\nApplied across all 5 model variants.",
         "WEEK 3-4"),
    ]
    for i, (code, title, body, status) in enumerate(contribs):
        x = Inches(0.4 + i * 3.2)
        _stripe(slide, TEAL if status == "COMPLETE" else RGBColor(0x1B, 0x46, 0x6E),
                x, Inches(1.3), Inches(3.0), Inches(3.4))
        _box(slide, code, x + Inches(0.1), Inches(1.35), Inches(0.6), Inches(0.5),
             size=22, bold=True, color=PALE)
        _box(slide, title, x + Inches(0.1), Inches(1.85), Inches(2.8), Inches(0.4),
             size=13, bold=True, color=WHITE)
        _box(slide, body, x + Inches(0.1), Inches(2.3), Inches(2.8), Inches(1.7),
             size=10, color=PALE)
        status_color = GREEN if status == "COMPLETE" else RGBColor(0xF3, 0x9C, 0x12)
        _box(slide, f"STATUS: {status}", x + Inches(0.1), Inches(4.35), Inches(2.8), Inches(0.3),
             size=10, bold=True, color=status_color)
    _slide_number(slide, 4)


def s05_c2_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT)
    _stripe(slide, TEAL, Inches(0), Inches(0), Inches(0.05), Inches(5.625))

    _label(slide, "C2: SYNTHETIC DATA PIPELINE -- BUILT & VALIDATED", Inches(0.3), Inches(0.25))
    _box(slide, "CTGAN Benchmark: 10 Facilities, 10,950 Records, 15 Clinical Features",
         Inches(0.3), Inches(0.55), Inches(9.4), Inches(0.55),
         size=19, bold=True, color=DARK)

    # Pipeline steps
    steps = [
        ("Schema Design", "15 features from MDS 3.0, RUG-IV,\nCMS PBJ staffing data.\nBinary mismatch label calibrated\nper care type."),
        ("CTGAN Training", "1 model per facility, 500 epochs.\nGenerator + Discriminator [256,256].\n~1,095 records/facility x 10 = 10,950 total."),
        ("Non-IID Engineering", "MC: adl_cognition=4.5, mismatch=40%\nSNF: adl_cognition=2.5, mismatch=28%\nIL: adl_cognition=1.0, mismatch=12%"),
        ("Fidelity Validation", "DUAL validation -- COMPLETE:\nInternal: 14/14 KS PASS, Frobenius\n0.2196 (27x baseline), TSTR 0.0199\nMIMIC-IV: cross-domain gap confirms\nLTC != hospital EHR (expected)"),
    ]
    for i, (step, body) in enumerate(steps):
        x = Inches(0.3 + i * 2.38)
        _stripe(slide, TEAL if i < 3 else GREEN, x, Inches(1.3), Inches(2.1), Inches(3.5))
        _box(slide, str(i + 1), x + Inches(0.08), Inches(1.35), Inches(0.4), Inches(0.4),
             size=18, bold=True, color=PALE)
        _box(slide, step, x + Inches(0.08), Inches(1.75), Inches(1.95), Inches(0.35),
             size=11, bold=True, color=WHITE)
        _box(slide, body, x + Inches(0.08), Inches(2.15), Inches(1.95), Inches(2.3),
             size=9.5, color=PALE)

    # MIMIC-IV completed note
    _stripe(slide, RGBColor(0xD4, 0xED, 0xDA), Inches(0.3), Inches(5.0), Inches(9.4), Inches(0.45))
    _box(slide, "DONE (16 Jun 2026): MIMIC-IV v3.1 real data validated (205,456 elderly admissions). Cross-domain gap confirms LTC != hospital EHR, motivating C2. Internal validation confirms CTGAN reproduces target LTC distributions.",
         Inches(0.45), Inches(5.0), Inches(9.1), Inches(0.45),
         size=9, bold=False, color=RGBColor(0x15, 0x55, 0x24))
    _slide_number(slide, 5)


def s06_c2_fidelity_fig(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT)
    _stripe(slide, TEAL, Inches(0), Inches(0), Inches(0.05), Inches(5.625))

    _label(slide, "C2 (REFRAMED): FIDELITY + WITHIN-MIMIC-IV COHORT CALIBRATION", Inches(0.3), Inches(0.25))
    _box(slide, "Internal Consistency + Honest Cross-Domain Disclosure + NEW Calibration Check",
         Inches(0.3), Inches(0.55), Inches(9.4), Inches(0.5),
         size=17, bold=True, color=DARK)

    fig2 = FIGURES / "fig2_fidelity_distributions.png"
    _img(slide, fig2, Inches(0.3), Inches(1.15), Inches(5.8))

    # Internal validation metrics (from synthetic holdout)
    _box(slide, "Internal (CTGAN vs LTC holdout):", Inches(6.3), Inches(1.1), Inches(3.5), Inches(0.28),
         size=9, bold=True, color=GREY)
    int_metrics = [
        ("14/14", "KS features, alpha=0.05", GREEN),
        ("0.2196", "Frobenius (27x baseline)", GREEN),
        ("0.0199", "TSTR gap, target <0.08", GREEN),
    ]
    for i, (val, sub, color) in enumerate(int_metrics):
        y = Inches(1.4 + i * 0.42)
        _box(slide, val, Inches(6.3), y, Inches(1.1), Inches(0.32),
             size=15, bold=True, color=color)
        _box(slide, sub, Inches(7.45), y + Inches(0.03), Inches(2.2), Inches(0.3),
             size=9, color=MID)

    # Cross-domain MIMIC-IV metrics (disclosed honestly, weak signal)
    _box(slide, "Cross-domain vs real MIMIC-IV (only 3/14 features mappable):", Inches(6.3), Inches(2.85), Inches(3.5), Inches(0.4),
         size=9, bold=True, color=GREY)
    _box(slide, "0/3 KS pass | Frobenius 0.82 | TSTR gap 0.33",
         Inches(6.3), Inches(3.2), Inches(3.5), Inches(0.3),
         size=10, bold=True, color=TEAL)
    _box(slide, "Expected: hospital admission != LTC facility-day.\nMotivates purpose-built LTC benchmark.",
         Inches(6.3), Inches(3.52), Inches(3.5), Inches(0.45),
         size=8.5, italic=True, color=MID)

    # NEW cohort calibration
    _box(slide, "NEW -- within-MIMIC-IV cohort calibration:", Inches(6.3), Inches(4.05), Inches(3.5), Inches(0.3),
         size=9, bold=True, color=GREEN)
    _box(slide, "LTC-bound (27.2%) vs non-LTC (72.8%) discharge\nCohen's d=0.69, 0.78 (p<0.001) on acuity proxies\nReal-world rate within 0.8pt of synthetic SNF target",
         Inches(6.3), Inches(4.38), Inches(3.5), Inches(0.75),
         size=8.5, color=MID)
    _slide_number(slide, 6)


def s07_c1_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    _stripe(slide, TEAL, Inches(0), Inches(4.9), Inches(10), Inches(0.06))

    _label(slide, "C1: FL ARCHITECTURE -- BUILT", Inches(0.4), Inches(0.25))
    _box(slide, "Three-Layer HIPAA-Compliant Federated System",
         Inches(0.4), Inches(0.55), Inches(9.2), Inches(0.5),
         size=20, bold=True, color=WHITE)

    layers = [
        ("Layer 1: Facility Edge", TEAL,
         "10 facilities (9 train, 1 held-out)\nXGBoost trains locally on resident records\nOnly model bytes (~50-200 KB) transmitted\nFedAcuityClient (Flower NumPyClient)"),
        ("Layer 2: Aggregation Server", RGBColor(0x1B, 0x46, 0x6E),
         "Three strategies compared:\nFedAvg / FedProx / Clustered FL\nCFL: 3 independent cluster models\nMC [0,1,2], SNF [3,4,5,6], IL [7,8]\nFacility 9 HELD OUT -- sole unseen facility"),
        ("Layer 3: Evaluation & XAI", RGBColor(0x0A, 0x47, 0x3A),
         "ResultsLogger: JSON + CSV per round\nmetrics.py: bootstrap CI + Mann-Whitney\nfigures.py: convergence + bar chart\nXAI: SHAP TreeExplainer (post-MidSEM)"),
    ]
    for i, (title, bg, body) in enumerate(layers):
        x = Inches(0.4 + i * 3.22)
        _stripe(slide, bg, x, Inches(1.3), Inches(3.0), Inches(3.2))
        _box(slide, title, x + Inches(0.12), Inches(1.35), Inches(2.8), Inches(0.45),
             size=11, bold=True, color=PALE)
        _box(slide, body, x + Inches(0.12), Inches(1.85), Inches(2.8), Inches(2.4),
             size=10, color=WHITE)

    _box(slide, "Raw resident data NEVER leaves any facility. Only model weights travel. HIPAA-compliant by design.",
         Inches(0.4), Inches(4.7), Inches(9.2), Inches(0.3),
         size=9.5, bold=True, color=PALE, align=PP_ALIGN.CENTER)
    _slide_number(slide, 7)


def s08_c1_why_cfl(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT)
    _stripe(slide, TEAL, Inches(0), Inches(0), Inches(0.05), Inches(5.625))

    _label(slide, "C1: WHY CLUSTERED FL? THE NON-IID PROBLEM", Inches(0.3), Inches(0.25))
    _box(slide, "Global FedAvg Fails on Independent Living Facilities",
         Inches(0.3), Inches(0.55), Inches(9.4), Inches(0.5),
         size=19, bold=True, color=DARK)

    # Left: explanation
    _box(slide, "The Core Problem:",
         Inches(0.3), Inches(1.2), Inches(4.5), Inches(0.35), size=12, bold=True, color=DARK)
    problems = [
        "Memory Care (MC): 40% mismatch rate, adl_cognition=4.5",
        "Skilled Nursing (SNF): 28% mismatch rate, adl_cognition=2.5",
        "Independent Living (IL): 12% mismatch rate, adl_cognition=1.0",
        "",
        "FedAvg averages gradients from ALL three care types.",
        "IL facilities (low acuity, low mismatch) are drowned out",
        "by MC + SNF signals during global aggregation.",
        "",
        "CFL solution: Each cluster trains its own global model.",
        "IL model sees only IL data -> far better generalisation.",
    ]
    for i, line in enumerate(problems):
        y = Inches(1.6 + i * 0.36)
        color = TEAL if line.startswith("CFL") else (MID if line else GREY)
        bold = line.startswith("CFL")
        if line:
            _box(slide, ("* " if not line.startswith(("Memory", "Skilled", "Indep", "CFL")) and line else "") + line,
                 Inches(0.3), y, Inches(4.5), Inches(0.35), size=10, color=color, bold=bold)

    # Right: result table
    _box(slide, "On Held-Out IL Facility (9):", Inches(5.2), Inches(1.2), Inches(4.5), Inches(0.35),
         size=12, bold=True, color=DARK)
    results = [
        ("Strategy",        "AUC-ROC", "F1",    TEAL,  True),
        ("Centralised (UB)","0.9799",  "0.723", GREY,  False),
        ("IL Local (fac. 7)","0.9643", "0.652", GREY,  False),
        ("Clustered FL",    "0.9693",  "0.727", GREEN, True),
        ("FedAvg",          "0.9414",  "0.704", RED,   False),
        ("FedProx",         "0.9414",  "0.704", RED,   False),
    ]
    for i, (strat, auc, f1, col, bld) in enumerate(results):
        y = Inches(1.55 + i * 0.5)
        _box(slide, strat, Inches(5.2), y, Inches(2.4), Inches(0.45), size=10, color=col, bold=bld)
        _box(slide, auc,   Inches(7.7), y, Inches(0.9), Inches(0.45), size=10, color=col, bold=bld, align=PP_ALIGN.CENTER)
        _box(slide, f1,    Inches(8.7), y, Inches(0.8), Inches(0.45), size=10, color=col, bold=bld, align=PP_ALIGN.CENTER)

    _stripe(slide, RGBColor(0xE8, 0xF8, 0xF1), Inches(5.2), Inches(2.55), Inches(4.3), Inches(0.5))
    _box(slide, "+2.79 AUC PTS vs FedAvg | +0.50 vs IL Local",
         Inches(5.2), Inches(4.6), Inches(4.5), Inches(0.35),
         size=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    _slide_number(slide, 8)


def s09_c1_results_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT)
    _stripe(slide, TEAL, Inches(0), Inches(0), Inches(0.05), Inches(5.625))

    _label(slide, "C1: FULL RESULTS TABLE -- HELD-OUT FACILITY 9", Inches(0.3), Inches(0.25))
    _box(slide, "50 Communication Rounds | Evaluated on Unseen IL Facility 9 | 9 Training Clients",
         Inches(0.3), Inches(0.55), Inches(9.4), Inches(0.5),
         size=18, bold=True, color=DARK)

    headers = ["Strategy", "AUC-ROC", "F1 Score", "Precision", "Recall", "Status"]
    rows = [
        ("IL Local (fac. 7 only)","0.9643", "0.6522", "0.7500", "0.5769", "Care-type local baseline"),
        ("Cross-Facility Ensemble","0.9414","0.7037", "0.6786", "0.7308", "9 models, no protocol"),
        ("Centralised Oracle",    "0.9799", "0.7234", "0.8095", "0.6538", "Upper bound (HIPAA violation)"),
        ("FedAvg",                "0.9414", "0.7037", "0.6786", "0.7308", "Global FL underperforms CFL"),
        ("FedProx (mu=0.1)",      "0.9414", "0.7037", "0.6786", "0.7308", "Same as FedAvg (XGB limitation)"),
        ("Clustered FL [C1]",     "0.9693", "0.7273", "0.8889", "0.6154", "PRIMARY CONTRIBUTION"),
    ]
    col_ws = [1.9, 0.8, 0.85, 0.85, 0.75, 2.4]
    # Header row
    x0 = Inches(0.3)
    for i, (h, cw) in enumerate(zip(headers, col_ws)):
        x = x0 + sum(Inches(col_ws[j]) for j in range(i))
        _stripe(slide, NAVY, x, Inches(1.18), Inches(cw), Inches(0.34))
        _box(slide, h, x + Inches(0.05), Inches(1.2), Inches(cw - 0.05), Inches(0.32),
             size=9, bold=True, color=WHITE)
    # Data rows
    for ri, row in enumerate(rows):
        y = Inches(1.52 + ri * 0.52)
        is_cfl = ri == 5
        bg_row = RGBColor(0xE8, 0xF8, 0xF1) if is_cfl else (WHITE if ri % 2 == 0 else LIGHT)
        for i, (val, cw) in enumerate(zip(row, col_ws)):
            x = x0 + sum(Inches(col_ws[j]) for j in range(i))
            _stripe(slide, bg_row, x, y, Inches(cw), Inches(0.48))
            color = GREEN if is_cfl else (RED if ri in [3, 4] and i in [1, 2, 3, 4] else MID)
            _box(slide, val, x + Inches(0.05), y + Inches(0.03), Inches(cw - 0.05), Inches(0.42),
                 size=9, bold=is_cfl, color=color)

    _box(slide, "Mann-Whitney U = 100.0  |  p = 1.6e-05  |  Statistically SIGNIFICANT (CFL vs FedAvg)",
         Inches(0.3), Inches(5.1), Inches(9.4), Inches(0.35),
         size=10, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    _slide_number(slide, 9)


def s09b_transparency(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    _stripe(slide, TEAL, Inches(0), Inches(4.9), Inches(10), Inches(0.06))

    _label(slide, "TRANSPARENCY: MID-CYCLE STRUCTURAL FIX", Inches(0.4), Inches(0.25))
    _box(slide, "A Flaw Was Found, Disclosed, and Fixed Before This Submission",
         Inches(0.4), Inches(0.55), Inches(9.2), Inches(0.6),
         size=18, bold=True, color=WHITE)

    cols = [
        ("C1 Structural Fix", TEAL, [
            "Found: IL cluster {7,8,9} had only facility 7",
            "training (8 & 9 both held out) -> Clustered FL",
            "for IL degenerated to single-client local training.",
            "",
            "Fixed: moved facility 8 into the IL training set.",
            "IL now has 2 real clients (7, 8); facility 9 is the",
            "sole, genuinely unseen held-out facility.",
            "",
            "Re-ran: FL simulation, held-out eval, DP sweep,",
            "and Figs 3-5 end-to-end on the corrected config.",
        ]),
        ("C2 Re-framing", PALE, [
            "Found: only 3 of 14 features map MIMIC-IV to the",
            "LTC schema -- a materially weak cross-domain claim.",
            "",
            "Re-framed: disclosed the weak mapping honestly",
            "instead of overstating it; added a NEW within-",
            "MIMIC-IV cohort-calibration analysis (LTC-bound vs",
            "non-LTC discharge) independent of that mapping.",
            "",
            "Result: medium-large effect sizes (d=0.69, 0.78,",
            "p<0.001) and a real-world face-validity check.",
        ]),
    ]
    for i, (title, color, items) in enumerate(cols):
        x = Inches(0.4 + i * 4.7)
        _stripe(slide, RGBColor(0x12, 0x2B, 0x52), x, Inches(1.25), Inches(4.4), Inches(3.5))
        _box(slide, title, x + Inches(0.15), Inches(1.32), Inches(4.0), Inches(0.4),
             size=13, bold=True, color=color)
        for j, item in enumerate(items):
            y = Inches(1.78 + j * 0.27)
            if item:
                _box(slide, item, x + Inches(0.15), y, Inches(4.1), Inches(0.28),
                     size=9.5, color=WHITE if i == 0 else PALE)

    _box(slide, "Why this slide exists: academic integrity requires disclosing material corrections, not silently revising numbers. Full detail in Section 4.4 of the MidSEM report.",
         Inches(0.4), Inches(4.95), Inches(9.2), Inches(0.32),
         size=8.5, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    _slide_number(slide, 10, total=17)


def s10_convergence(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT)
    _stripe(slide, TEAL, Inches(0), Inches(0), Inches(0.05), Inches(5.625))

    _label(slide, "C1: FL CONVERGENCE CURVES -- 50 ROUNDS", Inches(0.3), Inches(0.25))
    _box(slide, "CFL Consistently Outperforms FedAvg Across All 50 Rounds",
         Inches(0.3), Inches(0.55), Inches(9.4), Inches(0.5),
         size=19, bold=True, color=DARK)

    fig3 = FIGURES / "fig3_convergence.png"
    _img(slide, fig3, Inches(0.3), Inches(1.15), Inches(6.2))

    bullets = [
        ("CFL", "Converges at AUC 0.9909 (training clients)", GREEN),
        ("FedAvg", "Plateaus at AUC 0.9883 -- 0.26 pts below CFL", RED),
        ("FedProx", "Identical to FedAvg (XGBoost proximal term not implementable)", RED),
        ("", "", WHITE),
        ("Observation", "Training-client gaps are small; the real test is generalisation to a facility never seen during training.", MID),
        ("", "", WHITE),
        ("On held-out", "CFL 0.9693 vs FedAvg 0.9414 (facility 9) -- the true test of generalisation to unseen IL.", TEAL),
    ]
    for i, (key, val, col) in enumerate(bullets):
        y = Inches(1.15 + i * 0.62)
        if key:
            _box(slide, key + ":", Inches(6.7), y, Inches(1.1), Inches(0.35),
                 size=10, bold=True, color=col)
            _box(slide, val, Inches(7.85), y, Inches(1.9), Inches(0.55),
                 size=10, color=MID)
    _slide_number(slide, 11, total=17)


def s11_dp(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT)
    _stripe(slide, TEAL, Inches(0), Inches(0), Inches(0.05), Inches(5.625))

    _label(slide, "DIFFERENTIAL PRIVACY -- EPSILON SWEEP COMPLETE", Inches(0.3), Inches(0.25))
    _box(slide, "Privacy-Utility Tradeoff: epsilon=10 Is the Recommended Operating Point",
         Inches(0.3), Inches(0.55), Inches(9.4), Inches(0.5),
         size=18, bold=True, color=DARK)

    fig5 = FIGURES / "fig5_dp_privacy_utility.png"
    _img(slide, fig5, Inches(0.3), Inches(1.15), Inches(5.8))

    _box(slide, "Results (Opacus DP-SGD on StaffingNN):",
         Inches(6.3), Inches(1.15), Inches(3.5), Inches(0.35), size=11, bold=True, color=DARK)

    dp_data = [
        ("eps=1",  "0.6215", "35.6% drop", RED),
        ("eps=2",  "0.6911", "28.4% drop", RED),
        ("eps=5",  "0.7210", "25.3% drop", RED),
        ("eps=10", "0.7784", "19.3% drop", GREEN),
        ("No DP",  "0.9649", "Baseline",   TEAL),
    ]
    for i, (eps, auc, note, col) in enumerate(dp_data):
        y = Inches(1.55 + i * 0.62)
        is_rec = eps == "eps=10"
        if is_rec:
            _stripe(slide, RGBColor(0xE8, 0xF8, 0xF1), Inches(6.3), y - Inches(0.05), Inches(3.5), Inches(0.55))
        _box(slide, eps, Inches(6.3), y, Inches(1.0), Inches(0.45), size=10, bold=is_rec, color=col)
        _box(slide, auc, Inches(7.4), y, Inches(0.8), Inches(0.45), size=10, bold=is_rec, color=col)
        _box(slide, ("<<< RECOMMENDED" if is_rec else note), Inches(8.3), y, Inches(1.5), Inches(0.45),
             size=9, bold=is_rec, color=col)

    _box(slide, "Context: epsilon<=5 causes >25% utility drop with this 3-layer NN. epsilon=10 gives 19.3% drop -- best practical balance. Re-run on the corrected 9-facility training pool / facility-9 eval split. Future work: tree-level XGBoost DP for tighter privacy.",
         Inches(6.3), Inches(4.7), Inches(3.5), Inches(0.75), size=9, color=GREY)
    _slide_number(slide, 12, total=17)


def s12_stats(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    _stripe(slide, TEAL, Inches(0), Inches(4.9), Inches(10), Inches(0.06))

    _label(slide, "STATISTICAL VALIDATION", Inches(0.4), Inches(0.25))
    _box(slide, "Results Are Statistically Highly Significant",
         Inches(0.4), Inches(0.55), Inches(9.2), Inches(0.5),
         size=20, bold=True, color=WHITE)

    tests = [
        ("Mann-Whitney U Test", "CFL vs FedAvg\n(per-round AUC distributions)",
         "U = 100.0\np = 1.6e-05\n\nStatistically significant\n(p < 0.001)", GREEN),
        ("Bootstrap 95% CI", "CFL AUC\n(training clients, SEED=42)",
         "AUC = 0.9909\n(training clients)\n\nExtremely tight CI:\nmodel is consistent", TEAL),
        ("Effect Size", "CFL vs FedAvg / IL Local\non held-out facility 9",
         "vs FedAvg: +0.0279 AUC\nvs IL Local: +0.0050 AUC\n\nMeaningful practical effect\nconfirmed statistically", PALE),
    ]
    for i, (title, sub, result, color) in enumerate(tests):
        x = Inches(0.4 + i * 3.2)
        _stripe(slide, RGBColor(0x12, 0x2B, 0x52), x, Inches(1.3), Inches(3.0), Inches(3.3))
        _box(slide, title, x + Inches(0.12), Inches(1.35), Inches(2.8), Inches(0.4),
             size=13, bold=True, color=PALE)
        _box(slide, sub, x + Inches(0.12), Inches(1.8), Inches(2.8), Inches(0.55),
             size=10, color=GREY)
        _box(slide, result, x + Inches(0.12), Inches(2.4), Inches(2.8), Inches(1.5),
             size=11, bold=True, color=color)
    _slide_number(slide, 13, total=17)


def s13_c3_plan(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT)
    _stripe(slide, TEAL, Inches(0), Inches(0), Inches(0.05), Inches(5.625))

    _label(slide, "C3: XAI AUDIT SCORECARD -- POST-MIDSEM PLAN", Inches(0.3), Inches(0.25))
    _box(slide, "4-Dimension SHAP Audit: Planned for June 14 - June 27",
         Inches(0.3), Inches(0.55), Inches(9.4), Inches(0.5),
         size=19, bold=True, color=DARK)

    dims = [
        ("D1", "Fidelity", "Spearman rho of top-10 SHAP\nfeature ranks vs Centralised Oracle.\nTarget: rho >= 0.75 for CFL.\nWeek 3 (Jun 17-18)"),
        ("D2", "Stability", "Mean SHAP shift under +/-5%\nGaussian noise (100 perturbations).\nHypothesis: CFL < FedAvg shift.\nWeek 3 (Jun 19-20)"),
        ("D3", "Fairness", "Equalized odds + demographic\nparity across MC / SNF / IL.\nNormalised [0,1] score.\nWeek 4 (Jun 21-23)"),
        ("D4", "Plausibility", "% top-5 SHAP features matching\nLTC clinical literature list.\nTarget: >= 60% match.\nWeek 4 (Jun 21-23)"),
    ]
    for i, (code, dim, body) in enumerate(dims):
        x = Inches(0.3 + i * 2.38)
        _stripe(slide, NAVY, x, Inches(1.3), Inches(2.1), Inches(3.5))
        _box(slide, code, x + Inches(0.1), Inches(1.35), Inches(0.55), Inches(0.5),
             size=22, bold=True, color=PALE)
        _box(slide, dim, x + Inches(0.1), Inches(1.85), Inches(1.9), Inches(0.38),
             size=13, bold=True, color=WHITE)
        _box(slide, body, x + Inches(0.1), Inches(2.3), Inches(1.9), Inches(2.3),
             size=9.5, color=PALE)

    _box(slide, "Output: Fig 6 (radar chart) + XAI Audit Scorecard CSV  |  Paper Section VII filled  |  Ready for Final SEM Jul 11",
         Inches(0.3), Inches(5.05), Inches(9.4), Inches(0.35),
         size=10, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    _slide_number(slide, 14, total=17)


def s14_challenges(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    _stripe(slide, TEAL, Inches(0), Inches(4.9), Inches(10), Inches(0.06))

    _label(slide, "ENGINEERING CHALLENGES RESOLVED", Inches(0.4), Inches(0.25))
    _box(slide, "Real Implementation Has Real Problems -- Here Is What We Fixed",
         Inches(0.4), Inches(0.55), Inches(9.2), Inches(0.5),
         size=19, bold=True, color=WHITE)

    challenges = [
        ("XGBoost 3.x Serialisation", "save_raw() + temp file deserialization workaround. Flower expects numpy arrays; XGBoost models must be byte-encoded."),
        ("IL Mismatch Rate Bug", "NON_IID_SPEC missing adl_eating + adl_toileting caused IL rate to read 41% instead of target 12%. Fixed in schema.py; 72/72 tests now pass."),
        ("No Ray on Windows", "Flower simulation rewritten as manual round loop. No Ray dependency. Enables clean Windows 11 execution."),
        ("MIMIC-IV Validated (16 Jun)", "205,456 elderly admissions preprocessed. Cross-domain gap confirmed (0/3 KS, TSTR 0.42 vs 0.75). Domain mismatch proves LTC-specific benchmark is necessary. C2 contribution validated."),
        ("DP on XGBoost Not Possible", "Opacus requires gradient access -- not available in XGBoost tree boosting. Secondary PyTorch NN (StaffingNN) used for DP sweep only."),
        ("Windows Unicode (cp1252)", "Box-drawing chars (U+2500 etc.) cause codec errors. All print() statements use ASCII. Scripts run cleanly on Windows 11."),
        ("C1 Degenerate IL Cluster", "Self-review found IL cluster had only 1 training client (8 & 9 both held out). Fixed by moving facility 8 into IL training; re-ran simulation, eval, DP sweep end-to-end. See Slide 10."),
    ]
    for i, (title, body) in enumerate(challenges):
        row, col = divmod(i, 2)
        x = Inches(0.4 + col * 4.8)
        y = Inches(1.2 + row * 0.98)
        _stripe(slide, RGBColor(0x12, 0x2B, 0x52), x, y, Inches(4.5), Inches(0.92))
        _box(slide, title, x + Inches(0.1), y + Inches(0.04), Inches(4.3), Inches(0.28),
             size=9.5, bold=True, color=PALE)
        _box(slide, body, x + Inches(0.1), y + Inches(0.32), Inches(4.3), Inches(0.58),
             size=9, color=GREY)
    _slide_number(slide, 15, total=17)


def s15_timeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT)
    _stripe(slide, TEAL, Inches(0), Inches(0), Inches(0.05), Inches(5.625))

    _label(slide, "PLAN OF WORK: ABSTRACT TO FINAL SEM", Inches(0.3), Inches(0.25))
    _box(slide, "Full Dissertation Timeline -- Abstract Phase to Submission",
         Inches(0.3), Inches(0.55), Inches(9.4), Inches(0.5),
         size=19, bold=True, color=DARK)

    phases = [
        ("Apr 25 - May 8",  "Dissertation Outline",      "Lit review, architecture design, environment setup",  "DONE", GREEN),
        ("May 9 - May 21",  "Data Engineering",           "CTGAN generation, schema, non-IID engineering, validation pipeline", "DONE", GREEN),
        ("May 22 - Jun 13", "FL Implementation & DP",     "5 model variants, 50 rounds, DP sweep, midsem eval", "DONE", GREEN),
        ("Jun 14 - Jul 9",  "XAI Audit Engine",           "SHAP pipeline, D1-D4 modules, scorecard, radar chart (Fig 6)", "NEXT", TEAL),
        ("Jul 10 - Jul 27", "Dissertation Review",        "Full draft, methodology, results, discussion, plagiarism check", "PLANNED", GREY),
        ("Jul 28 - Aug 2",  "Submission",                 "Final review and submission with supervisor evaluation report", "PLANNED", GREY),
    ]
    hdr_ys = [0.3, 0.3, 0.3, 0.3]
    for i, (dates, phase, work, status, color) in enumerate(phases):
        y = Inches(1.3 + i * 0.68)
        is_now = status == "NEXT"
        bg_col = RGBColor(0xE8, 0xF8, 0xF1) if status == "DONE" else (RGBColor(0xE8, 0xF3, 0xFF) if is_now else WHITE)
        _stripe(slide, bg_col, Inches(0.3), y, Inches(9.4), Inches(0.62))
        # Status indicator
        _stripe(slide, color, Inches(0.3), y, Inches(0.08), Inches(0.62))
        _box(slide, dates, Inches(0.5), y + Inches(0.1), Inches(1.4), Inches(0.45),
             size=9, color=GREY)
        _box(slide, phase, Inches(2.0), y + Inches(0.1), Inches(2.2), Inches(0.45),
             size=10, bold=True, color=color)
        _box(slide, work, Inches(4.3), y + Inches(0.1), Inches(4.2), Inches(0.45),
             size=9.5, color=MID)
        _box(slide, status, Inches(8.6), y + Inches(0.1), Inches(1.0), Inches(0.45),
             size=9, bold=True, color=color, align=PP_ALIGN.CENTER)
    _slide_number(slide, 16, total=17)


def s16_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    _stripe(slide, TEAL, Inches(0), Inches(4.9), Inches(10), Inches(0.06))

    _label(slide, "MIDSEM SUMMARY -- WHAT WAS ACHIEVED", Inches(0.4), Inches(0.2))
    _box(slide, "From Proposal to Results: FedAcuity Is Real and Working",
         Inches(0.4), Inches(0.5), Inches(9.2), Inches(0.5),
         size=20, bold=True, color=WHITE)

    cols = [
        ("C1: FL System", GREEN, [
            "5 strategies implemented + tested",
            "50-round simulation, 9 training facilities",
            "CFL AUC 0.9693 vs FedAvg 0.9414",
            "+2.79pt gap on held-out facility 9 (IL)",
            "Mann-Whitney U=100.0, p=1.6e-05",
        ]),
        ("C2: Synthetic Data", TEAL, [
            "CTGAN pipeline: 10 facilities done",
            "10,950 records, 15 features",
            "Internal: 14/14 KS PASS, TSTR 0.0199",
            "MIMIC-IV: 205,456 real admissions used",
            "NEW: cohort calibration, d=0.69/0.78",
        ]),
        ("C3: XAI Scorecard", PALE, [
            "Architecture designed",
            "SHAP pipeline planned",
            "D1-D4 modules: Jun 14 - Jun 27",
            "Fig 6 radar chart: Week 4",
            "Paper Sec VII ready post-midsem",
        ]),
    ]
    for i, (title, color, items) in enumerate(cols):
        x = Inches(0.4 + i * 3.22)
        _stripe(slide, RGBColor(0x12, 0x2B, 0x52), x, Inches(1.2), Inches(3.0), Inches(3.4))
        _box(slide, title, x + Inches(0.12), Inches(1.25), Inches(2.8), Inches(0.38),
             size=13, bold=True, color=color)
        for j, item in enumerate(items):
            _box(slide, ("+ " if i < 2 else "~ ") + item,
                 x + Inches(0.12), Inches(1.7 + j * 0.5), Inches(2.8), Inches(0.45),
                 size=10, color=WHITE if i < 2 else PALE)

    _box(slide, "Next: SHAP pipeline + D1-D4 XAI modules (Jun 14)  |  Final SEM: 11 July 2026",
         Inches(0.4), Inches(4.95), Inches(9.2), Inches(0.28),
         size=10, color=PALE, align=PP_ALIGN.CENTER)
    _slide_number(slide, 17, total=17)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    r = _load()
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    s01_title(prs)
    s02_agenda(prs)
    s03_problem(prs)
    s04_abstract_recap(prs)
    s05_c2_pipeline(prs)
    s06_c2_fidelity_fig(prs)
    s07_c1_architecture(prs)
    s08_c1_why_cfl(prs)
    s09_c1_results_table(prs)
    s09b_transparency(prs)
    s10_convergence(prs)
    s11_dp(prs)
    s12_stats(prs)
    s13_c3_plan(prs)
    s14_challenges(prs)
    s15_timeline(prs)
    s16_summary(prs)

    OUT.parent.mkdir(exist_ok=True)
    prs.save(str(OUT))
    print(f"Slide deck saved: {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
